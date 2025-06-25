import os
import matplotlib.pyplot as plt
import scipy.stats
from itertools import combinations
import numpy as np
import pandas as pd
import traceback

from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
import tempfile


def write_summary_stats(df, filename, master_file='experiment_summary.xlsx'):
    """
    Write summary statistics to a master Excel file, with one line per experimental group.
    Control values are included as separate columns next to treatment values.
    
    Parameters:
    df (pandas.DataFrame): DataFrame containing the experimental data
    filename (str): Name of the source file being analyzed
    master_file (str): Path to the master Excel file
    """
    import pandas as pd
    import numpy as np
    import scipy.stats
    import os
    from openpyxl import load_workbook
    
    # Define columns to include
    base_columns = ['Date', 'Filename', 'Hair', 'Treatment', 'Treatment time', 'Name']
    metric_columns = ['BREAK STRESS']
    
    # Define friendly names mapping
    metric_friendly_names = {
        'BREAK STRESS': 'Break Stress',
    }
    
    # Function to check if a name contains control-related words
    def is_control_group(name):
        control_terms = ['control', 'controls', 'ctrl', 'base', 'wash']
        return any(term in name.lower() for term in control_terms)
    
    # Get all unique groups and separate control from treatment groups
    all_groups = df['Name'].unique()
    control_groups = [group for group in all_groups if is_control_group(group)]
    treatment_groups = [group for group in all_groups if not is_control_group(group)]
    
    # Create summary DataFrame
    summary_data = []
    
    # Process each treatment group
    for group in treatment_groups:
        group_data = df[df['Name'] == group]
        
        # Get the first occurrence of base columns and add filename
        row_data = {col: group_data[col].iloc[0] if col != 'Filename' else filename 
                   for col in base_columns}
        
        # Initialize hair count variable
        hair_count = []
        
        # Initialize control group name variable
        control_group_used = None
        
        # Calculate stats for each metric
        for metric in metric_columns:
            if metric in df.columns:
                # Remove outliers for treatment group
                cleaned_data, _ = remove_outliers(group_data[metric])
                hair_count.append(len(cleaned_data))
                
                # Calculate mean and std for treatment group
                treatment_mean = round(cleaned_data.mean(), 2)
                treatment_std = round(cleaned_data.std(), 2)
                
                # Get friendly name for the metric
                friendly_name = metric_friendly_names[metric]
                
                # Find corresponding control group and calculate its stats
                control_mean = None
                control_std = None
                control_count = None
                p_value = None
                
                if control_groups:
                    # Use the first control group found (you might want to refine this logic)
                    control_group = control_groups[0]
                    control_group_used = control_group  # Store the control group name
                    control_data = df[df['Name'] == control_group][metric]
                    control_cleaned, _ = remove_outliers(control_data)
                    
                    if len(control_cleaned) > 0:
                        control_mean = round(control_cleaned.mean(), 2)
                        control_std = round(control_cleaned.std(), 2)
                        control_count = len(control_cleaned)
                        
                        # Perform statistical test between treatment and control
                        if len(cleaned_data) > 0:
                            try:
                                # Perform Mann-Whitney U test
                                stat, p_value = scipy.stats.mannwhitneyu(cleaned_data, control_cleaned, 
                                                                 alternative='two-sided')
                                # Round p-value to 4 decimal places
                                p_value = round(p_value, 4)
                            except Exception as e:
                                print(f"Statistical test failed for {group} vs {control_group}: {e}")
                                p_value = None
                
                # Store treatment values
                row_data[f'{friendly_name} Mean'] = treatment_mean
                row_data[f'{friendly_name} Std'] = treatment_std
                row_data[f'{friendly_name} Count'] = len(cleaned_data)
                
                # Store control values in adjacent columns
                row_data[f'{friendly_name} Control Mean'] = control_mean
                row_data[f'{friendly_name} Control Std'] = control_std
                row_data[f'{friendly_name} Control Count'] = control_count
                
                # Store p-value
                row_data[f'{friendly_name} p-value'] = p_value
        
        # Add control group name
        row_data['Control Group Used'] = control_group_used
        
        # Add hair count as the last column
        if hair_count:
            row_data['Mean # of hairs after outliers removal (across metrics)'] = round(np.mean(hair_count), 1)
        else:
            row_data['Mean # of hairs after outliers removal (across metrics)'] = 0
            
        summary_data.append(row_data)

    
    # Create DataFrame from summary data
    new_summary_df = pd.DataFrame(summary_data)
    
    # Load existing file and append new data
    if os.path.exists(master_file):
        try:
            # Read existing data
            existing_df = pd.read_excel(master_file)
            
            # Create unique identifiers for comparison
            # Use a more robust comparison that handles potential data type issues
            def create_identifier(row):
                return f"{str(row['Name'])}_{str(row['Date'])}_{str(row['Filename'])}"
            
            existing_identifiers = set(existing_df.apply(create_identifier, axis=1))
            new_identifiers = new_summary_df.apply(create_identifier, axis=1)
            
            # Filter out entries that already exist
            mask = ~new_identifiers.isin(existing_identifiers)
            new_entries_to_add = new_summary_df[mask]
            
            if len(new_entries_to_add) > 0:
                print(f"Adding {len(new_entries_to_add)} new entries to {master_file}")
                # Use pd.concat instead of append (which is deprecated)
                combined_df = pd.concat([existing_df, new_entries_to_add], ignore_index=True)
            else:
                print(f"No new entries to add for {filename}")
                combined_df = existing_df
                
        except Exception as e:
            print(f"Error reading existing file {master_file}: {e}")
            print("Creating new file with current data")
            combined_df = new_summary_df
    else:
        print(f"Creating new master file: {master_file}")
        combined_df = new_summary_df
    
    try:
        # Write to Excel with error handling
        combined_df.to_excel(master_file, index=False, engine='openpyxl')
        print(f"Successfully updated {master_file}")
    except PermissionError:
        print(f"Permission denied writing to {master_file}. Make sure the file is not open in Excel.")
        # Try with a backup filename
        backup_file = master_file.replace('.xlsx', '_backup.xlsx')
        combined_df.to_excel(backup_file, index=False, engine='openpyxl')
        print(f"Data saved to backup file: {backup_file}")
    except Exception as e:
        print(f"Error writing to Excel file: {e}")
        # Save as CSV as fallback
        csv_file = master_file.replace('.xlsx', '.csv')
        combined_df.to_csv(csv_file, index=False)
        print(f"Data saved as CSV: {csv_file}")
def reorder_by_names(df, name_list):
    """
    Groups rows by name and orders groups according to name_list
    """
    mask = pd.CategoricalIndex(df['Name'], 
                              categories=name_list, 
                              ordered=True)
    return df.set_index(mask).sort_index().reset_index(drop=True)

def add_missing_records(df):
    """
    Add rows for missing record numbers in the RECORD column.
    The function assumes RECORD column should contain consecutive integers from 1 to max(RECORD).
    
    Parameters:
    df (pandas.DataFrame): DataFrame containing a 'RECORD' column
    
    Returns:
    pandas.DataFrame: DataFrame with added rows for missing record numbers
    """
    # Ensure RECORD column exists
    if 'RECORD' not in df.columns:
        raise ValueError("DataFrame must contain a 'RECORD' column")
    
    # Get the current max record number
    max_record = int(df['RECORD'].max())
    
    # Create a set of existing record numbers
    existing_records = set(df['RECORD'])
    existing_records = [int(x) for x in existing_records]
    
    # Find missing record numbers
    missing_records = [x for x in range(1, max_record + 1) if x not in existing_records]
    print(f"Missing records: {missing_records}")
    
    # If no missing records, return original dataframe
    if not missing_records:
        return df
    
    # Create new rows for missing records
    # Fill other columns with NaN/empty values
    new_rows = pd.DataFrame({
        'RECORD': missing_records,
        **{col: [np.nan] * len(missing_records) for col in df.columns if col != 'RECORD'}
    })
    
    # Concatenate original df with new rows
    result_df = pd.concat([df, new_rows], ignore_index=True)
    
    # Sort by RECORD number
    result_df = result_df.sort_values('RECORD').reset_index(drop=True)
    all_records = list(result_df['RECORD'])
    print(f"Records order check ", all_records[:6])
    
    return result_df

def switch_units_triple(df):
    # multiple all the values in the column
    # BREAK STRESS by 10^6/101.971621297792 / comes from what Devora was doing in her file
    # note from June 23, 2025: it's to make units in Newton the weird 101.97 ...
    df['BREAK STRESS'] = df['BREAK STRESS'] * 10**6/101.971621297792
    # and EMOD to gigapascal
    df['ELASTIC EMOD'] = df['ELASTIC EMOD'] * 1e-9
    # and thouness to mega joules
    df['TOUGHNESS'] = df['TOUGHNESS'] * 1e-6
    return df

def switch_units_single(df):
    df['EMOD(*)(#)'] = df['EMOD(*)(#)'] * 1e-9
    return df

def get_df_from_file(filepath, skip):
    try:
        # Read the entire file at once instead of line by line
        with open(filepath, 'r', encoding='utf-8') as f:
            # Convert to list to avoid iterator issues
            all_lines = list(f)
            
        # Skip the header lines
        data_lines = all_lines[skip:]
        
        # Process the data lines
        raw_data = []
        for line in data_lines:
            fields = line.strip().split('\t')
            # Pad with None if needed
            fields.extend([None] * (15 - len(fields)))
            raw_data.append(fields)

        # Convert to DataFrame
        df = pd.DataFrame(raw_data)
        # Replace empty values and common NA indicators
        df = df.replace(['', 'NA', 'null'], np.nan)
        
        return df
        
    except Exception as e:
        print(f"Error reading file {filepath}: {e}")
        print(f"Full traceback: {traceback.format_exc()}")
        return None

def replace_invalid_floats(df):
    """
    Replace '1.#IOe+000' and '-1.#IOe+000' values with numpy.nan in a DataFrame.
    Works on both string representations and actual float values.
    
    Parameters:
    -----------
    df : pandas.DataFrame
        The input DataFrame containing invalid float values
        
    Returns:
    --------
    pandas.DataFrame
        A new DataFrame with invalid float values replaced by nan
    """
    # Create a copy to avoid modifying the original DataFrame
    df_cleaned = df.copy()
    
    # Replace string representations
    df_cleaned = df_cleaned.replace(['1.#IOe+000', '-1.#IOe+000'], np.nan)
    
    # For each numeric column, replace invalid float values
    numeric_columns = df_cleaned.select_dtypes(include=[np.number]).columns
    for col in numeric_columns:
        # Check for invalid float values
        mask = df_cleaned[col].apply(lambda x: 
            isinstance(x, float) and (str(x) == '1.#IOe+000' or str(x) == '-1.#IOe+000'))
        # Replace with nan where mask is True
        df_cleaned.loc[mask, col] = np.nan
    
    return df_cleaned

def clean_single(df, max_100=True):
    # replace invalid floats
    df_clean = replace_invalid_floats(df)
    # set the header to the rows with RECORD
    df_clean.columns = df_clean.iloc[0]
    df_clean = df_clean[1:]
    df_clean.reset_index(drop=True, inplace=True)
    # removes row with the unit
    df_clean = df_clean.drop(index=0)
    df_clean.reset_index(drop=True, inplace=True)
    # remove empty rows at the end
    mask = pd.to_numeric(df_clean['RECORD'], errors='coerce').notnull()
    df_clean = df_clean[mask]
    df_clean.reset_index(drop=True, inplace=True)
    # save all the values as float 
    df_clean = df_clean.astype(float)
    # swtich units
    df_clean = switch_units_single(df_clean)
    # add missing records
    df_clean = add_missing_records(df_clean)
    
    if max_100:
        # if there is a value 101 in the colum RECORD, delete the row
        df_clean = df_clean[df_clean['RECORD'] != 101]
        
    return df_clean

def clean_triple(df, max_100=True):
    # replace invalid floats
    df_clean = replace_invalid_floats(df)
    # set the header to the rows with RECORD
    df_clean.columns = df_clean.iloc[0]
    df_clean = df_clean[1:]
    df_clean.reset_index(drop=True, inplace=True)
    # remove the first row
    df_clean = df_clean.drop(index=0)
    df_clean.reset_index(drop=True, inplace=True)
    # removes row with the unit
    df_clean = df_clean.drop(index=0)
    df_clean.reset_index(drop=True, inplace=True)
    # remove empty rows at the end
    mask = pd.to_numeric(df_clean['RECORD'], errors='coerce').notnull()
    df_clean = df_clean[mask]
    df_clean.reset_index(drop=True, inplace=True)
    # save all the values as float 
    df_clean = df_clean.astype(float)
    # switch units for break stress
    df_clean = switch_units_triple(df_clean)
    # add missing records
    df_clean = add_missing_records(df_clean)

    if max_100:
        # if there is a value 101 in the colum RECORD, delete the row
        df_clean = df_clean[df_clean['RECORD'] != 101]

    return df_clean

def format_number(number):
    """Format number: use scientific notation if more than 4 digits"""
    if abs(number) >= 10000:
        return f"{number:.2e}"
    return f"{number:.2f}"

def remove_outliers(data):
    Q1 = np.percentile(data, 25)
    Q3 = np.percentile(data, 75)
    IQR = Q3 - Q1
    lower_bound = Q1 - 1.5 * IQR
    upper_bound = Q3 + 1.5 * IQR

    to_keep = data[(data >= lower_bound) & (data <= upper_bound)]
    removed = data[(data < lower_bound) | (data > upper_bound)]
    assert len(to_keep) + len(removed) == len(data)

    return to_keep, removed

def create_scatter_plot(df, x_col, y_col, save_dir=None, filename=None, figsize=(10, 6), show=False, verbose=False):
    """
    Create a scatter plot with correlation coefficient for two columns in a dataframe
    and optionally save it to a specified directory.
    
    Parameters:
    -----------
    df : pandas.DataFrame
        The input dataframe containing the data
    x_col : str
        Name of the column to plot on x-axis
    y_col : str
        Name of the column to plot on y-axis
    save_dir : str, optional
        Directory path where to save the plot. If None, plot is only displayed
    filename : str, optional
        Name of the file to save the plot. If None and save_dir is specified,
        a default name will be generated
    figsize : tuple, optional
        Size of the figure (width, height) in inches
        
    Returns:
    --------
    None (displays and optionally saves the plot)
    """
    # Calculate correlation
    correlation = df[x_col].corr(df[y_col])
    
    # Create the scatter plot
    plt.figure(figsize=figsize)
    ax = plt.gca()
    
    # Remove top and right spines
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    
    # Create the scatter plot
    plt.scatter(df[x_col], df[y_col], alpha=0.5)
    
    # Add labels and title
    plt.xlabel(x_col)
    plt.ylabel(y_col)
    plt.title(f'{y_col} vs {x_col}')
    
    # Add correlation text
    plt.text(0.02, 0.98, f'Correlation: {correlation:.3f}', 
             transform=plt.gca().transAxes,
             verticalalignment='top',
             bbox=dict(boxstyle='round', facecolor='white', alpha=0.8))
    
    # Add grid for better readability
    plt.grid(True, linestyle='--', alpha=0.7)
    
    # Adjust layout to prevent text cutoff
    plt.tight_layout()
    
    # Save the plot if directory is specified
    if save_dir:
        # Create directory if it doesn't exist
        os.makedirs(save_dir, exist_ok=True)
        
        # Generate default filename if none provided
        if filename is None:
            filename = f'scatter_{x_col}_vs_{y_col}.png'
        
        # Ensure filename has an extension
        if not filename.endswith(('.png', '.jpg', '.jpeg', '.pdf')):
            filename += '.png'
        
        # Create full path
        filepath = os.path.join(save_dir, filename)
        
        # Save the plot
        plt.savefig(filepath, bbox_inches='tight', dpi=300)
        if verbose:
            print(f"Plot saved to: {filepath}")
    
    # Show the plot
    if show:
        plt.show()
        plt.close()

def create_boxplot(df, metric_column, ymin, ymax, experiment_name, group_column='Name', 
                   figsize=(10, 6), verbose=False):
    
    """Create a compact boxplot with statistical test, median values, mean ± std, vertically stacked significance indicators,
    and individual data points (outliers removed)"""
    
    print(f"\nCreating boxplot for {metric_column}")
    # Create plot
    fig, ax = plt.subplots(figsize=figsize)
    
    # Define colors for control and non-control groups
    control_color = '#E0E0E0'  # light grey
    experiment_color = '#B2DFDB'  # light teal
    
    # Function to check if a name contains control-related words
    def is_control_group(name):
        control_terms = ['control', 'controls', 'ctrl']
        return any(term in name.lower() for term in control_terms)
    
    # Create lists to store cleaned data for boxplot
    cleaned_data = []
    group_names = df[group_column].unique()
    
    all_removed = []
    for group in group_names:
        group_data = df[df[group_column] == group][metric_column]
        cleaned_group_data, removed = remove_outliers(group_data)
        print("Removed outliers for", group, ":", removed)
        all_removed.append(removed)
        cleaned_data.append(cleaned_group_data)
    # merge the dataframe
    all_removed = pd.concat(all_removed)
    
    # Create boxplot with cleaned data
    bp = ax.boxplot(cleaned_data,
                    patch_artist=True,
                    medianprops={'color': 'red', 'linewidth': 1.5},
                    boxprops={'color': 'black'},
                    whiskerprops={'color': 'black'},
                    capprops={'color': 'black'},
                    showfliers=False)
    
    # Color the boxes based on whether they're control groups
    for patch, group_name in zip(bp['boxes'], group_names):
        color = control_color if is_control_group(group_name) else experiment_color
        patch.set_facecolor(color)
        patch.set_alpha(0.7)
    
    # Add individual points with jitter (using cleaned data)
    for i, (group_name, group_data) in enumerate(zip(group_names, cleaned_data)):
        color = control_color if is_control_group(group_name) else experiment_color
        
        # Create jitter
        x = np.random.normal(i + 1, 0.04, size=len(group_data))
        
        # Plot points with black edges
        ax.scatter(x, group_data, 
                  color=color,
                  edgecolor='black',
                  linewidth=0.5,
                  alpha=0.5,
                  s=20,
                  zorder=2)
    
    # Rest of the function remains the same
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.set_xticklabels(group_names, rotation=45, ha='right')
    
    def group_markers(markers_list):
        """Group identical markers together and organize different markers in rows"""
        if not markers_list:
            return []
            
        # Sort markers to group identical ones together
        sorted_markers = sorted(markers_list, key=lambda x: x[0])
        
        # Group identical markers
        grouped = []
        current_group = [sorted_markers[0]]
        
        for marker, color in sorted_markers[1:]:
            if marker == current_group[0][0]:
                current_group.append((marker, color))
            else:
                grouped.append(current_group)
                current_group = [(marker, color)]
        grouped.append(current_group)
        
        return grouped
    
    # Dictionary to store significance markers and their colors for each group
    significance_dict = {group: [] for group in group_names}
    
    # List of base markers to use for different pairs
    base_markers = ['†', '‡', '§', '¶', '#', '@', '©', '®', '°', '∞', '∆', '∑', '√', '∫']
    marker_idx = 0
    
    # Define colors for different pairs (using a colorful qualitative colormap)
    pair_colors = plt.cm.Dark2(np.linspace(0, 1, len(list(combinations(group_names, 2)))))
    
    # Perform pairwise Mann-Whitney U tests
    pairs = list(combinations(group_names, 2))
    
    for pair_idx, (group1_name, group2_name) in enumerate(pairs):
        idx1 = list(group_names).index(group1_name)
        idx2 = list(group_names).index(group2_name)
        
        group1_data = cleaned_data[idx1]
        group2_data = cleaned_data[idx2]
        
        stat, p_value = scipy.stats.mannwhitneyu(group1_data, group2_data, alternative='two-sided')
        
        if p_value < 0.05:
            current_marker = base_markers[marker_idx % len(base_markers)]
            marker_idx += 1
            
            current_color = pair_colors[pair_idx]
            
            if p_value < 0.001:
                marker_with_significance = current_marker * 3
            elif p_value < 0.01:
                marker_with_significance = current_marker * 2
            else:
                marker_with_significance = current_marker
                
            significance_dict[group1_name].append((marker_with_significance, current_color))
            significance_dict[group2_name].append((marker_with_significance, current_color))
            
            print(f"Stats. significance: {group1_name} and {group2_name}: p = {p_value:.4f}")
    
    # Add median values, means, and significance markers
    for i, group in enumerate(group_names):
        median = np.median(cleaned_data[i])
        mean = np.mean(cleaned_data[i])
        std = np.std(cleaned_data[i])
        
        # Get the actual y-axis limits
        y_min, y_max = ax.get_ylim()
        y_range = y_max - y_min
        
        # Add median value
        ax.text(i+1, y_max + y_range*0.02, f'M: {format_number(median)}',
                horizontalalignment='center', fontsize=8)
        
        # Add mean ± std
        ax.text(i+1, y_max + y_range*0.08, f'μ: {format_number(mean)}±{format_number(std)}',
                horizontalalignment='center', fontsize=8)
        
        # Add significance markers if they exist
        if significance_dict[group]:
            grouped_markers = group_markers(significance_dict[group])
            
            for row_idx, marker_group in enumerate(grouped_markers):
                marker_text = ''.join(marker for marker, _ in marker_group)
                color = marker_group[0][1]
                
                # Calculate vertical position for significance markers
                vertical_offset = 0.14 + (row_idx * 0.05)
                
                ax.text(i+1, y_max + y_range*vertical_offset, marker_text,
                    horizontalalignment='center', fontsize=12, color=color)

    # After all annotations are added, adjust the plot limits
    max_rows = max(len(group_markers(markers)) if markers else 0 
                for markers in significance_dict.values())
    final_offset = 0.14 + (max_rows * 0.05) + 0.05  # Add a small padding at the top

    ax.set_ylim(y_min, y_max + y_range*final_offset)
    
    plt.suptitle('')
    plt.title('')
    plt.xlabel(f'{experiment_name}', fontsize=10)
    plt.ylabel(metric_column.replace('_', ' ').title(), fontsize=10)
    ax.yaxis.grid(True, linestyle='--', alpha=0.3)
    
    # Add legend for significance levels
    legend_text = '1 symbol: p < 0.05\n2 symbols: p < 0.01\n3 symbols: p < 0.001'
    plt.figtext(0.99, 0.01, legend_text, 
                horizontalalignment='right', 
                verticalalignment='bottom',
                fontsize=6,
                bbox=dict(facecolor='white', alpha=0.8, edgecolor='none'))

    if ymin is not None and ymax is not None:
        ax.set_ylim(ymin, ymax)
    
    plt.tight_layout()
    
    return fig, ax, all_removed

########################
# pdf / ppt functions

def add_to_powerpoint(fig, title, pptx_path, experiment_name=None, analysis_date=None):
    """
    Add or update a matplotlib figure in a PowerPoint presentation.
    If a slide with the same title exists, it will be updated rather than creating a new slide.
    
    Parameters:
    -----------
    fig : matplotlib.figure.Figure
        The figure to add to the presentation
    title : str
        Base title for the slide (e.g., "Break Stress Analysis")
    pptx_path : str
        Path to the PowerPoint file
    experiment_name : str, optional
        Name of the experiment or treatment group
    analysis_date : str
        Date of the analysis in any format
    """
    if analysis_date is None:
        raise ValueError("analysis_date must be provided")
        
    # Generate a unique slide title
    full_title = title
    if experiment_name:
        full_title = f"{title} - {experiment_name}"
    full_title = f"{full_title} ({analysis_date})"
    
    # Save the figure to a temporary file
    temp_img_path = tempfile.NamedTemporaryFile(suffix='.png', delete=False).name
    fig.savefig(temp_img_path, dpi=300, bbox_inches='tight')
    
    try:
        if os.path.exists(pptx_path):
            prs = Presentation(pptx_path)
            
            # Check each slide for matching title
            for slide in prs.slides:
                if hasattr(slide.shapes, 'title') and slide.shapes.title:
                    if slide.shapes.title.text == full_title:
                        # If matching slide found, delete all its shapes
                        shape_ids = [shape.shape_id for shape in slide.shapes]
                        for shape_id in shape_ids:
                            shape = slide.shapes._spTree.find(f'.//*[@id="{shape_id}"]')
                            if shape is not None:
                                slide.shapes._spTree.remove(shape)
                        
                        # Add new title and image
                        title_box = slide.shapes.title
                        title_box.text = full_title
                        title_box.text_frame.paragraphs[0].font.size = Pt(14)
                        
                        left = Inches(1)
                        top = Inches(1.5)
                        slide.shapes.add_picture(temp_img_path, left, top, height=Inches(5))
                        
                        prs.save(pptx_path)
                        return
        else:
            prs = Presentation()
            
        # If we get here, either the file is new or no matching slide was found
        # Use layout 5 (Title and Content)
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # Ensure title is added
        if slide.shapes.title:
            slide.shapes.title.text = full_title
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(14)
        
        # Add the image
        left = Inches(1)
        top = Inches(1.5)
        slide.shapes.add_picture(temp_img_path, left, top, height=Inches(5))
        
        # Save the presentation
        prs.save(pptx_path)
        
    finally:
        # Clean up temporary file
        if os.path.exists(temp_img_path):
            os.remove(temp_img_path)